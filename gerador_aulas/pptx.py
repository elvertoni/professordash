"""
Extrator de PPTX/PPT usando python-pptx.
Extrai: título, corpo, tabelas, notas do apresentador.
"""

import logging
from .base import ExtratorBase, ResultadoExtracao, TipoArquivo, limpar_texto

logger = logging.getLogger(__name__)


def extrair_pptx(arquivo) -> ResultadoExtracao:
    """
    Extrai todo o conteúdo de uma apresentação PowerPoint.
    Inclui texto de cada shape, tabelas e notas do apresentador.
    """
    try:
        from pptx import Presentation
        from pptx.util import Pt
    except ImportError:
        return ResultadoExtracao(
            conteudo='',
            tipo=TipoArquivo.PPTX,
            erro='python-pptx não instalado. Execute: pip install python-pptx',
        )

    try:
        prs = Presentation(arquivo)
        slides_texto = []

        for i, slide in enumerate(prs.slides, start=1):
            partes = [f'--- Slide {i} ---']

            # Ordena shapes por posição vertical (de cima para baixo)
            shapes_ordenados = sorted(
                slide.shapes,
                key=lambda s: (s.top or 0, s.left or 0)
            )

            for shape in shapes_ordenados:
                # Texto comum
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        texto = para.text.strip()
                        if texto:
                            # Detecta se é título pelo tamanho da fonte ou nome do placeholder
                            eh_titulo = (
                                shape.name.lower().startswith('title') or
                                _tem_fonte_grande(para)
                            )
                            prefixo = '## ' if eh_titulo else ''
                            partes.append(f'{prefixo}{texto}')

                # Tabelas
                if shape.has_table:
                    texto_tabela = _extrair_tabela_pptx(shape.table)
                    if texto_tabela:
                        partes.append(texto_tabela)

            # Notas do apresentador (muito úteis para enriquecer a aula)
            if slide.has_notes_slide:
                nota = slide.notes_slide.notes_text_frame.text.strip()
                if nota and nota != 'Clique para editar as notas do orador':
                    partes.append(f'\n[Nota do apresentador: {nota}]')

            if len(partes) > 1:  # ignora slides completamente vazios
                slides_texto.append('\n'.join(partes))

        conteudo = '\n\n'.join(slides_texto)
        conteudo = limpar_texto(conteudo)

        return ResultadoExtracao(
            conteudo=conteudo,
            tipo=TipoArquivo.PPTX,
            slides=len(prs.slides),
            metadados={
                'total_slides': len(prs.slides),
                'slides_com_conteudo': len(slides_texto),
            },
        )

    except Exception as e:
        logger.error(f'Erro ao extrair PPTX: {e}')
        return ResultadoExtracao(
            conteudo='',
            tipo=TipoArquivo.PPTX,
            erro=str(e),
        )


def _extrair_tabela_pptx(tabela) -> str:
    """Converte tabela do PPTX para texto com separadores."""
    linhas = []
    for linha in tabela.rows:
        celulas = [cell.text.strip() for cell in linha.cells]
        linhas.append(' | '.join(celulas))
    return '\n'.join(linhas)


def _tem_fonte_grande(paragrafo, limiar_pt: int = 24) -> bool:
    """Verifica se o parágrafo tem fonte grande (possível título)."""
    for run in paragrafo.runs:
        if run.font.size and run.font.size.pt >= limiar_pt:
            return True
    return False
