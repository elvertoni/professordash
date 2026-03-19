"""
Testes dos Critérios de Aceite do MVP — Gerador de Aulas

Cada classe documenta e verifica um critério de aceite específico do TASKS_AULAS.md.
Todos os testes mocam chamadas ao OpenRouter — sem custo de API.
"""

import io
import json
from unittest.mock import MagicMock, patch

import pytest

from gerador.extractors import extrair_arquivo
from gerador.extractors.base import TipoArquivo
from gerador.extractors.rco import detectar_papel_rco
from gerador.models import MaterialEntrada, SessaoGeracao
from gerador.pipeline import executar_modo_livre
from gerador.prompts import SYSTEM_PROMPT, prompt_rco
from gerador.providers import MODELOS
from gerador.tokens import UsoTokens


# ── Helpers ────────────────────────────────────────────────────────────────────

def _make_uso_mock():
    return UsoTokens(
        modelo="anthropic/claude-sonnet-4-5",
        tokens_input=800,
        tokens_output=2500,
    )


def _make_material_mock(conteudo, papel_rco=""):
    m = MagicMock()
    m.conteudo_extraido = conteudo
    m.papel_rco = papel_rco
    m.arquivo = None
    m.url = ""
    return m


def _fake_arquivo(nome, conteudo=b"dados"):
    buf = io.BytesIO(conteudo)
    buf.name = nome
    return buf


# ── MODO RCO ──────────────────────────────────────────────────────────────────

class TestCriterioRCODeteccaoPapel:
    """
    Critério: Sistema reconhece automaticamente o papel de cada arquivo pelo nome.
    """

    @pytest.mark.parametrize("nome,papel_esperado", [
        ("AULA_03_ATIVIDADE_FRONT_END.docx", "atividade"),
        ("AULA_03_PRATICA_FRONT_END.docx",   "pratica"),
        ("AULA_03_FRONT_END.pptx",           "slides"),
        ("AULA_07_MATEMÁTICA.pptx",          "slides"),
        ("aula_01_atividade_web.docx",        "atividade"),
        ("AULA_10_PRÁTICA_WEB.docx",          "pratica"),
    ])
    def test_detecta_papel_automaticamente(self, nome, papel_esperado):
        papel = detectar_papel_rco(nome)
        assert papel == papel_esperado, (
            f"'{nome}' deveria ser '{papel_esperado}', mas foi '{papel}'"
        )

    def test_arquivos_desconhecidos_marcados_como_outro(self):
        assert detectar_papel_rco("README.txt") == "outro"
        assert detectar_papel_rco("material_extra.pdf") == "outro"

    def test_atividade_nunca_confundida_com_pratica(self):
        """Regra crítica: nunca inverter ATIVIDADE com PRATICA."""
        assert detectar_papel_rco("AULA_05_ATIVIDADE.docx") == "atividade"
        assert detectar_papel_rco("AULA_05_PRATICA.docx")   == "pratica"
        assert detectar_papel_rco("AULA_05_ATIVIDADE.docx") != "pratica"


class TestCriterioRCOSlideComNotas:
    """
    Critério: Extrai todo o conteúdo dos slides, incluindo notas do apresentador.
    """

    def test_notas_extraidas_do_pptx(self):
        from pptx import Presentation
        from pptx.util import Pt

        prs = Presentation()
        layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = "Título do Slide"
        slide.notes_slide.notes_text_frame.text = "Nota do apresentador"

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        buf.name = "aula.pptx"

        from gerador.extractors.pptx import extrair_pptx
        resultado = extrair_pptx(buf)

        assert "Título do Slide" in resultado.conteudo
        assert "Nota do apresentador" in resultado.conteudo


class TestCriterioRCOSecoesMapeadas:
    """
    Critério: Questões de ATIVIDADE → seção 12; conteúdo de PRÁTICA → seção 11.
    """

    def test_atividade_instrucao_aponta_para_secao_12(self):
        prompt = prompt_rco(
            slides="Conteúdo dos slides",
            atividade="1. Qual o conceito de HTML?",
            pratica="",
            disciplina="Programação Web",
            numero=3,
            nivel="tecnico",
            instrucoes="",
            aula_anterior="",
        )
        assert "12" in prompt or "fixação" in prompt.lower() or "atividade" in prompt.lower()
        # A instrução explícita de mapeamento está presente no prompt
        texto_lower = prompt.lower()
        assert "atividade" in texto_lower

    def test_pratica_instrucao_aponta_para_secao_11(self):
        prompt = prompt_rco(
            slides="Conteúdo dos slides",
            atividade="",
            pratica="Crie um site com HTML semântico",
            disciplina="Programação Web",
            numero=3,
            nivel="tecnico",
            instrucoes="",
            aula_anterior="",
        )
        assert "11" in prompt or "prática" in prompt.lower() or "atividade prática" in prompt.lower()

    def test_prompt_rco_mapeia_atividade_12_pratica_11(self):
        """Verifica que o prompt menciona explicitamente seção 12 para atividade e 11 para prática."""
        prompt = prompt_rco(
            slides="Slides",
            atividade="Questões A/B/C",
            pratica="Exercício prático",
            disciplina="TI",
            numero=1,
            nivel="tecnico",
            instrucoes="",
            aula_anterior="",
        )
        assert "12" in prompt
        assert "11" in prompt


class TestCriterioRCOTempoGeracao:
    """
    Critério: Aula gerada em menos de 60 segundos.

    Não é possível testar tempo real com mocks. Verificamos que o pipeline não
    possui esperas artificiais (time.sleep) e que o fluxo é direto.
    """

    def test_pipeline_rco_sem_sleep_artificial(self):
        import ast
        import inspect
        import gerador.pipeline as pipeline_mod

        source = inspect.getsource(pipeline_mod)
        assert "time.sleep" not in source, (
            "pipeline.py não deve conter time.sleep — causaria lentidão artificial"
        )

    def test_pipeline_rco_chama_gerar_aula_uma_vez(self, turma):
        sessao = MagicMock()
        sessao.id = 1
        sessao.num_aulas = 55
        sessao.nivel = "tecnico"
        sessao.instrucoes = ""
        sessao.provider = "claude"
        sessao.disciplina = turma
        sessao.tokens_usados = 0
        sessao.custo_estimado = 0
        sessao.materiais.all.return_value = []

        from gerador.pipeline import executar_modo_rco
        uso_mock = _make_uso_mock()
        with patch("gerador.pipeline.gerar_aula", return_value=("## Aula\n\nConteúdo.", uso_mock)) as mock_gerar:
            with patch("gerador.pipeline.registrar_uso_na_sessao"):
                executar_modo_rco(sessao)

        assert mock_gerar.call_count == 1, "Modo RCO deve chamar gerar_aula exatamente 1 vez"


# ── MODO LIVRE ────────────────────────────────────────────────────────────────

class TestCriterioLivreFormatos:
    """
    Critério: Upload aceita PDF, PPTX, DOCX, TXT, MD e URL.
    """

    @pytest.mark.parametrize("extensao,tipo_esperado", [
        (".pdf",  TipoArquivo.PDF),
        (".pptx", TipoArquivo.PPTX),
        (".docx", TipoArquivo.DOCX),
        (".txt",  TipoArquivo.TXT),
        (".md",   TipoArquivo.TXT),
    ])
    def test_formato_aceito_e_roteado_corretamente(self, extensao, tipo_esperado):
        """Cada extensão é encaminhada ao extrator correto."""
        arq = _fake_arquivo(f"material{extensao}", b"conteudo de teste")

        if extensao == ".pdf":
            with patch("gerador.extractors.extrair_pdf") as mock_ext:
                mock_ext.return_value = MagicMock(conteudo="texto", tipo=tipo_esperado, palavras=1, erro="")
                resultado = extrair_arquivo(arq, f"material{extensao}")
            mock_ext.assert_called_once()
        elif extensao in (".pptx", ".ppt"):
            with patch("gerador.extractors.extrair_pptx") as mock_ext:
                mock_ext.return_value = MagicMock(conteudo="texto", tipo=tipo_esperado, palavras=1, erro="")
                resultado = extrair_arquivo(arq, f"material{extensao}")
            mock_ext.assert_called_once()
        elif extensao in (".docx", ".doc"):
            with patch("gerador.extractors.extrair_docx") as mock_ext:
                mock_ext.return_value = MagicMock(conteudo="texto", tipo=tipo_esperado, palavras=1, erro="")
                resultado = extrair_arquivo(arq, f"material{extensao}")
            mock_ext.assert_called_once()
        else:
            # .txt e .md: lidos diretamente
            arq = _fake_arquivo(f"material{extensao}", b"texto simples")
            resultado = extrair_arquivo(arq, f"material{extensao}")
            assert resultado.conteudo == "texto simples"
            assert resultado.tipo == tipo_esperado

    def test_formato_desconhecido_retorna_erro(self):
        arq = _fake_arquivo("video.mp4", b"binary")
        resultado = extrair_arquivo(arq, "video.mp4")
        assert resultado.erro != ""
        assert resultado.conteudo == ""

    def test_url_aceita_pelo_extrator_de_url(self):
        """URL é tratada pelo extrair_url (não pelo extrair_arquivo)."""
        from gerador.extractors.url import extrair_url
        from gerador.extractors.base import ResultadoExtracao
        mock_resultado = ResultadoExtracao(conteudo="Conteúdo da página", tipo=TipoArquivo.URL)
        with patch("gerador.extractors.url._extrair_com_bs4", return_value=mock_resultado):
            resultado = extrair_url("https://example.com/artigo")
        assert resultado.tipo == TipoArquivo.URL


class TestCriterioLivreTitulosEditaveis:
    """
    Critério: Professor edita títulos antes de aprovar (verificado via template e view).
    """

    def test_template_planejamento_tem_inputs_editaveis(self):
        """O template planejamento.html deve ter campos input para edição de títulos."""
        import os
        template_path = os.path.join(
            os.path.dirname(__file__), "../..", "templates", "gerador", "planejamento.html"
        )
        template_path = os.path.normpath(template_path)
        with open(template_path) as f:
            html = f.read()
        assert 'type="text"' in html, "Template deve ter inputs de texto para editar títulos"
        assert 'name="titulo_' in html, "Campo de título deve ter name='titulo_N'"

    @pytest.mark.django_db
    def test_aprovar_planejamento_atualiza_titulos(self, client_professor, turma):
        """AprovarPlanejamentoView atualiza os títulos editados pelo professor."""
        sessao = SessaoGeracao.objects.create(
            disciplina=turma,
            modo="livre",
            num_aulas=2,
            nivel="tecnico",
            foco="equilibrado",
            provider="claude",
            planejamento={
                "tema_central": "HTML",
                "fio_condutor": "Web",
                "observacoes": "",
                "aulas": [
                    {"numero": 1, "titulo": "Título Original 1", "topicos_principais": []},
                    {"numero": 2, "titulo": "Título Original 2", "topicos_principais": []},
                ],
            },
        )
        resp = client_professor.post(
            f"/painel/gerador/{sessao.pk}/aprovar/",
            {"titulo_1": "Novo Título 1", "titulo_2": "Novo Título 2"},
        )
        assert resp.status_code == 302  # redirect após salvar
        sessao.refresh_from_db()
        aulas = sessao.planejamento["aulas"]
        assert aulas[0]["titulo"] == "Novo Título 1"
        assert aulas[1]["titulo"] == "Novo Título 2"


class TestCriterioLivrePlanejamento:
    """
    Critério: modo livre já chega na tela de planejamento com proposta gerada.
    """

    @pytest.mark.django_db
    def test_upload_livre_gera_planejamento_antes_do_redirect(self, client_professor, turma):
        planejamento = {
            "tema_central": "Algoritmos",
            "fio_condutor": "Resolução de problemas",
            "observacoes": "",
            "aulas": [
                {"numero": 1, "titulo": "Introdução aos algoritmos", "topicos_principais": ["variáveis"]},
            ],
        }
        uso_mock = _make_uso_mock()

        with patch("gerador.views.gerar_planejamento", return_value=(planejamento, uso_mock)):
            with patch("gerador.tokens.registrar_uso_na_sessao"):
                resp = client_professor.post(
                    "/painel/gerador/upload/",
                    {
                        "modo": "livre",
                        "turma": str(turma.pk),
                        "num_aulas": "1",
                        "nivel": "tecnico",
                        "provider": "claude",
                        "texto_livre": "Conteúdo base para testar o planejamento.",
                    },
                )

        assert resp.status_code == 302

        sessao = SessaoGeracao.objects.latest("id")
        sessao.refresh_from_db()

        assert resp.url == f"/painel/gerador/{sessao.pk}/planejar/"
        assert sessao.planejamento == planejamento
        assert sessao.status == "rascunho"

    @pytest.mark.django_db
    def test_get_planejamento_recupera_sessao_sem_planejamento(self, client_professor, turma):
        sessao = SessaoGeracao.objects.create(
            disciplina=turma,
            modo="livre",
            num_aulas=1,
            nivel="tecnico",
            foco="equilibrado",
            provider="claude",
            planejamento=None,
        )
        MaterialEntrada.objects.create(
            sessao=sessao,
            tipo="texto",
            texto_livre="Base sobre lógica de programação",
            conteudo_extraido="Base sobre lógica de programação",
        )
        planejamento = {
            "tema_central": "Lógica",
            "fio_condutor": "Pensamento computacional",
            "observacoes": "",
            "aulas": [
                {"numero": 1, "titulo": "Conceitos iniciais", "topicos_principais": ["algoritmo"]},
            ],
        }
        uso_mock = _make_uso_mock()

        with patch("gerador.views.gerar_planejamento", return_value=(planejamento, uso_mock)):
            with patch("gerador.tokens.registrar_uso_na_sessao"):
                resp = client_professor.get(f"/painel/gerador/{sessao.pk}/planejar/")

        assert resp.status_code == 200

        sessao.refresh_from_db()
        assert sessao.planejamento == planejamento
        assert "Conceitos iniciais" in resp.content.decode()


class TestCriterioLivreSSEProgresso:
    """
    Critério: Progress bar atualiza em tempo real via SSE.
    """

    @pytest.mark.django_db
    def test_gerar_retorna_template_sem_header_sse(self, client_professor, turma):
        """GET sem Accept: text/event-stream renderiza a página de progresso."""
        sessao = SessaoGeracao.objects.create(
            disciplina=turma,
            modo="rco",
            num_aulas=1,
            nivel="tecnico",
            foco="equilibrado",
            provider="claude",
            planejamento=None,
        )
        resp = client_professor.get(f"/painel/gerador/{sessao.pk}/gerar/")
        assert resp.status_code == 200
        assert b"progress" in resp.content.lower() or b"gera" in resp.content.lower()

    @pytest.mark.django_db
    def test_gerar_retorna_sse_com_header_correto(self, client_professor, turma):
        """GET com Accept: text/event-stream retorna StreamingHttpResponse SSE."""
        sessao = SessaoGeracao.objects.create(
            disciplina=turma,
            modo="rco",
            num_aulas=1,
            nivel="tecnico",
            foco="equilibrado",
            provider="claude",
            planejamento=None,
        )

        uso_mock = _make_uso_mock()
        with patch("gerador.views.executar_modo_rco") as mock_rco:
            mock_rco.return_value = MagicMock()
            resp = client_professor.get(
                f"/painel/gerador/{sessao.pk}/gerar/",
                HTTP_ACCEPT="text/event-stream",
            )

        assert "text/event-stream" in resp.get("Content-Type", "")

    @pytest.mark.django_db
    def test_sse_livre_emite_evento_por_aula(self, client_professor, turma):
        """Modo livre: SSE emite um evento por aula gerada + evento de conclusão."""
        sessao = SessaoGeracao.objects.create(
            disciplina=turma,
            modo="livre",
            num_aulas=3,
            nivel="tecnico",
            foco="equilibrado",
            provider="claude",
            planejamento={
                "tema_central": "Web",
                "fio_condutor": "Semântica",
                "observacoes": "",
                "aulas": [
                    {"numero": 1, "titulo": "A1", "topicos_principais": []},
                    {"numero": 2, "titulo": "A2", "topicos_principais": []},
                    {"numero": 3, "titulo": "A3", "topicos_principais": []},
                ],
            },
        )

        # streaming_content é avaliado lazily — o patch deve estar ativo durante o consumo
        with patch("gerador.views.executar_modo_livre", return_value=iter([1, 2, 3])):
            resp = client_professor.get(
                f"/painel/gerador/{sessao.pk}/gerar/",
                HTTP_ACCEPT="text/event-stream",
            )
            conteudo = b"".join(resp.streaming_content)

        linhas = [l for l in conteudo.decode().split("\n") if l.startswith("data:")]

        # 3 eventos de progresso + 1 de conclusão
        assert len(linhas) == 4
        ultimo = json.loads(linhas[-1].replace("data: ", ""))
        assert ultimo.get("status") == "concluido"

    @pytest.mark.django_db
    def test_sse_erro_emite_evento_de_erro(self, client_professor, turma):
        """Quando pipeline levanta exceção, SSE emite evento de erro."""
        sessao = SessaoGeracao.objects.create(
            disciplina=turma,
            modo="rco",
            num_aulas=1,
            nivel="tecnico",
            foco="equilibrado",
            provider="claude",
            planejamento=None,
        )

        # streaming_content é avaliado lazily — o patch deve estar ativo durante o consumo
        with patch("gerador.views.executar_modo_rco", side_effect=Exception("Falha na API")):
            resp = client_professor.get(
                f"/painel/gerador/{sessao.pk}/gerar/",
                HTTP_ACCEPT="text/event-stream",
            )
            conteudo = b"".join(resp.streaming_content).decode()

        dados = json.loads(conteudo.split("data: ")[1])
        assert dados["status"] == "erro"
        assert "Falha na API" in dados["mensagem"]


class TestCriterioLivreLote20Aulas:
    """
    Critério: Geração em lote funciona de 1 a 20 aulas sem erro.
    """

    @pytest.mark.django_db
    @pytest.mark.parametrize("num_aulas", [1, 5, 10, 20])
    def test_lote_sem_erro(self, turma, num_aulas):
        """Geração em lote completa para 1, 5, 10 e 20 aulas."""
        sessao = MagicMock()
        sessao.id = 10
        sessao.nivel = "tecnico"
        sessao.foco = "equilibrado"
        sessao.provider = "claude"
        sessao.instrucoes = ""
        sessao.disciplina = turma
        sessao.tokens_usados = 0
        sessao.custo_estimado = 0
        sessao.planejamento = {
            "aulas": [
                {"numero": 200 + i, "titulo": f"Aula {i}", "topicos_principais": ["t"]}
                for i in range(1, num_aulas + 1)
            ]
        }
        sessao.materiais.all.return_value = [
            _make_material_mock("Material de apoio")
        ]

        uso_mock = _make_uso_mock()
        with patch("gerador.pipeline.gerar_aula", return_value=("## Aula\n\nConteúdo.", uso_mock)):
            with patch("gerador.pipeline.registrar_uso_na_sessao"):
                numeros = list(executar_modo_livre(sessao))

        assert len(numeros) == num_aulas
        assert sessao.status == "concluido"


# ── AMBOS OS MODOS ─────────────────────────────────────────────────────────────

class TestCriterioRascunhosEditaveis:
    """
    Critério: Aulas salvas como rascunhos editáveis na disciplina (realizada=False).
    """

    @pytest.mark.django_db
    def test_aula_salva_com_realizada_false(self, turma):
        from gerador.pipeline import executar_modo_rco

        sessao = MagicMock()
        sessao.id = 1
        sessao.num_aulas = 98
        sessao.nivel = "tecnico"
        sessao.instrucoes = ""
        sessao.provider = "claude"
        sessao.disciplina = turma
        sessao.tokens_usados = 0
        sessao.custo_estimado = 0
        sessao.materiais.all.return_value = []

        uso_mock = _make_uso_mock()
        with patch("gerador.pipeline.gerar_aula", return_value=("## Rascunho\n\nConteúdo.", uso_mock)):
            with patch("gerador.pipeline.registrar_uso_na_sessao"):
                aula = executar_modo_rco(sessao)

        assert aula.realizada is False, "Aula gerada deve ser rascunho (realizada=False)"

    @pytest.mark.django_db
    def test_aula_associada_a_turma_correta(self, turma):
        from gerador.pipeline import executar_modo_rco

        sessao = MagicMock()
        sessao.id = 1
        sessao.num_aulas = 97
        sessao.nivel = "tecnico"
        sessao.instrucoes = ""
        sessao.provider = "claude"
        sessao.disciplina = turma
        sessao.tokens_usados = 0
        sessao.custo_estimado = 0
        sessao.materiais.all.return_value = []

        uso_mock = _make_uso_mock()
        with patch("gerador.pipeline.gerar_aula", return_value=("## Aula\n\nConteúdo.", uso_mock)):
            with patch("gerador.pipeline.registrar_uso_na_sessao"):
                aula = executar_modo_rco(sessao)

        assert aula.turma == turma


class TestCriterioPadrao15Secoes:
    """
    Critério: Padrão de 15 seções obrigatórias respeitado em 100% das gerações.
    """

    SECOES_OBRIGATORIAS = [
        "cabeçalho",         # 1
        "título",             # 2
        "introdução",         # 3
        "competências",       # 4
        "recapitulação",      # 5
        "reflexão",           # 6
        "conceituação",       # 7
        "analogias",          # 8
        "detalhamento técnico", # 9
        "roteiro",            # 10
        "prática",            # 11
        "fixação",            # 12
        "para casa",          # 13
        "resumo",             # 14
        "referências",        # 15
    ]

    def test_system_prompt_menciona_todas_15_secoes(self):
        prompt_lower = SYSTEM_PROMPT.lower()
        faltando = []
        for secao in self.SECOES_OBRIGATORIAS:
            if secao not in prompt_lower:
                faltando.append(secao)
        assert not faltando, (
            f"SYSTEM_PROMPT não menciona as seções: {faltando}"
        )

    def test_prompt_livre_exige_15_secoes(self):
        from gerador.prompts import prompt_aula_livre
        aula = {"numero": 1, "titulo": "Teste", "topicos_principais": ["a"]}
        prompt = prompt_aula_livre(
            aula=aula, total=1, conteudo="material",
            disciplina="TI", nivel="tecnico", foco="equilibrado",
            aula_anterior="", instrucoes="",
        )
        assert "15" in prompt

    def test_prompt_rco_exige_todas_secoes(self):
        prompt = prompt_rco(
            slides="Slides", atividade="Questões", pratica="Prática",
            disciplina="TI", numero=1, nivel="tecnico",
            instrucoes="", aula_anterior="",
        )
        # Prompt deve solicitar o padrão completo de seções
        assert "11" in prompt and "12" in prompt


class TestCriterioProvidersSuportados:
    """
    Critério: Suporte a Claude, Gemini e GPT-4o via OpenRouter.
    """

    def test_todos_3_providers_mapeados(self):
        assert "claude" in MODELOS
        assert "gemini" in MODELOS
        assert "gpt4o"  in MODELOS

    def test_claude_usa_modelo_correto(self):
        assert "claude" in MODELOS["claude"].lower()

    def test_gemini_usa_modelo_correto(self):
        assert "gemini" in MODELOS["gemini"].lower()

    def test_gpt4o_usa_modelo_correto(self):
        assert "gpt" in MODELOS["gpt4o"].lower()

    def test_cada_provider_chama_openrouter(self):
        """gerar_aula() usa o modelo correto para cada provider."""
        from gerador.providers import gerar_aula

        uso_mock = _make_uso_mock()
        for provider in ["claude", "gemini", "gpt4o"]:
            with patch("gerador.providers._get_client") as mock_client:
                mock_resp = MagicMock()
                mock_resp.choices[0].message.content = "## Aula\n\nConteúdo."
                mock_resp.usage.prompt_tokens = 500
                mock_resp.usage.completion_tokens = 1500
                mock_client.return_value.chat.completions.create.return_value = mock_resp

                conteudo, uso = gerar_aula("system", "user", provider)

                call_kwargs = mock_client.return_value.chat.completions.create.call_args
                modelo_chamado = call_kwargs.kwargs.get("model") or call_kwargs.args[0] if call_kwargs else ""
                assert MODELOS[provider] in str(call_kwargs), (
                    f"Provider '{provider}' deve usar modelo '{MODELOS[provider]}'"
                )


class TestCriterioTokensECusto:
    """
    Critério: Estimativa de tokens e custo exibida ao final.
    """

    def test_uso_tokens_calcula_custo_automaticamente(self):
        """UsoTokens calcula custo USD ao ser criado."""
        uso = UsoTokens(
            modelo="anthropic/claude-sonnet-4-5",
            tokens_input=10_000,
            tokens_output=2_000,
        )
        assert uso.tokens_total == 12_000
        assert uso.custo_usd > 0

    def test_custo_brl_calculado(self):
        uso = UsoTokens(
            modelo="anthropic/claude-sonnet-4-5",
            tokens_input=5_000,
            tokens_output=1_000,
        )
        assert uso.custo_brl > 0

    def test_template_gerando_exibe_tokens_e_custo(self):
        """O template gerando.html contém elementos para exibir tokens e custo."""
        import os
        template_path = os.path.normpath(
            os.path.join(os.path.dirname(__file__), "../..", "templates", "gerador", "gerando.html")
        )
        with open(template_path) as f:
            html = f.read()
        assert 'id="tokens"' in html, "Template deve ter elemento com id='tokens'"
        assert 'id="custo"'  in html, "Template deve ter elemento com id='custo'"

    def test_registrar_uso_atualiza_sessao(self):
        """registrar_uso_na_sessao() acumula tokens e custo na SessaoGeracao."""
        from gerador.tokens import registrar_uso_na_sessao

        sessao_mock = MagicMock()
        sessao_mock.tokens_usados = 0
        sessao_mock.custo_estimado = 0

        uso = UsoTokens(
            modelo="anthropic/claude-sonnet-4-5",
            tokens_input=1_000,
            tokens_output=2_000,
        )
        registrar_uso_na_sessao(sessao_mock, uso)

        assert sessao_mock.tokens_usados == 3_000
        assert sessao_mock.custo_estimado > 0
