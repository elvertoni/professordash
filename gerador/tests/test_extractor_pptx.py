"""
Testes unitários para gerador/extractors/pptx.py
Cria apresentações em memória com python-pptx para testar a extração real.
"""

import io

import pytest
from pptx import Presentation
from pptx.util import Inches, Pt

from gerador.extractors.base import TipoArquivo
from gerador.extractors.pptx import extrair_pptx


# ── Fixtures ──────────────────────────────────────────────────────────────────

def _make_pptx(slides_data: list[dict]) -> io.BytesIO:
    """
    Cria PPTX em memória.
    slides_data: lista de dicts com 'titulo', 'corpo', 'nota' (opcionais).
    """
    prs = Presentation()
    layout = prs.slide_layouts[1]  # título + conteúdo

    for dados in slides_data:
        slide = prs.slides.add_slide(layout)

        if "titulo" in dados:
            slide.shapes.title.text = dados["titulo"]

        if "corpo" in dados:
            corpo = slide.placeholders[1]
            corpo.text = dados["corpo"]

        if "nota" in dados:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = dados["nota"]

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    buf.name = "aula.pptx"
    return buf


# ── Testes ────────────────────────────────────────────────────────────────────

class TestExtrairPptx:

    def test_extrai_titulo_e_corpo(self):
        """Deve extrair título e corpo de cada slide."""
        pptx = _make_pptx([
            {"titulo": "Introdução ao HTML", "corpo": "HTML é a linguagem da web."},
        ])
        resultado = extrair_pptx(pptx)

        assert resultado.ok
        assert resultado.tipo == TipoArquivo.PPTX
        assert "Introdução ao HTML" in resultado.conteudo
        assert "HTML é a linguagem da web." in resultado.conteudo

    def test_extrai_multiplos_slides(self):
        """Deve extrair conteúdo de todos os slides."""
        pptx = _make_pptx([
            {"titulo": "Slide 1", "corpo": "Conteúdo 1"},
            {"titulo": "Slide 2", "corpo": "Conteúdo 2"},
            {"titulo": "Slide 3", "corpo": "Conteúdo 3"},
        ])
        resultado = extrair_pptx(pptx)

        assert resultado.slides == 3
        assert "Slide 1" in resultado.conteudo
        assert "Slide 2" in resultado.conteudo
        assert "Slide 3" in resultado.conteudo

    def test_extrai_notas_do_apresentador(self):
        """Notas do apresentador devem ser incluídas no conteúdo."""
        pptx = _make_pptx([
            {
                "titulo": "CSS Flexbox",
                "corpo": "Alinhamento de elementos",
                "nota": "Lembrar de mostrar o exemplo ao vivo no CodePen.",
            }
        ])
        resultado = extrair_pptx(pptx)

        assert "Lembrar de mostrar o exemplo" in resultado.conteudo
        assert "[Nota do apresentador:" in resultado.conteudo

    def test_separadores_de_slide_presentes(self):
        """Cada slide deve ter o separador '--- Slide N ---'."""
        pptx = _make_pptx([
            {"titulo": "A", "corpo": "X"},
            {"titulo": "B", "corpo": "Y"},
        ])
        resultado = extrair_pptx(pptx)

        assert "--- Slide 1 ---" in resultado.conteudo
        assert "--- Slide 2 ---" in resultado.conteudo

    def test_retorna_erro_para_arquivo_invalido(self):
        """Arquivo que não é PPTX deve retornar ResultadoExtracao com erro."""
        arquivo_invalido = io.BytesIO(b"isso nao e um pptx valido")
        arquivo_invalido.name = "invalido.pptx"
        resultado = extrair_pptx(arquivo_invalido)

        assert not resultado.ok
        assert resultado.erro != ""

    def test_contagem_de_palavras_preenchida(self):
        """ResultadoExtracao deve calcular palavras automaticamente."""
        pptx = _make_pptx([{"titulo": "Tema", "corpo": "palavra1 palavra2 palavra3"}])
        resultado = extrair_pptx(pptx)

        assert resultado.palavras > 0
