"""
Testes unitários para gerador/extractors/rco.py

Cobre:
- detectar_papel_rco(): identificação por nome de arquivo
- montar_conjunto_de_uploads(): montagem a partir de request.FILES
- extrair_rco(): orquestração da extração
- montar_conteudo_rco(): consolidação por papel
- Regra crítica: ATIVIDADE → seção 12, PRATICA → seção 11 (nunca inverter)
"""

import io
from unittest.mock import MagicMock, patch

import pytest
from docx import Document
from pptx import Presentation

from gerador.extractors.base import ResultadoExtracao, TipoArquivo
from gerador.extractors.rco import (
    ConjuntoRCO,
    PapelRCO,
    detectar_numero_aula,
    detectar_papel_rco,
    extrair_rco,
    montar_conteudo_rco,
    montar_conjunto_de_uploads,
)


# ── Fixtures ──────────────────────────────────────────────────────────────────

def _make_pptx_buffer(titulo: str = "Slide de teste") -> io.BytesIO:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = titulo
    slide.placeholders[1].text = "Conteúdo do slide de teste"
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


def _make_docx_buffer(texto: str = "Questão de teste") -> io.BytesIO:
    doc = Document()
    doc.add_paragraph(texto)
    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf


# ── Testes: detectar_papel_rco ────────────────────────────────────────────────

class TestDetectarPapelRco:

    @pytest.mark.parametrize("nome,esperado", [
        # Atividade — diversas variações de nome
        ("AULA_03_ATIVIDADE_FRONT_END.docx",    PapelRCO.ATIVIDADE),
        ("aula_03_atividade_front_end.docx",    PapelRCO.ATIVIDADE),
        ("AULA_10_ATIVIDADE_PROGRAMACAO.docx",  PapelRCO.ATIVIDADE),
        # Prática — com e sem acento
        ("AULA_03_PRATICA_FRONT_END.docx",      PapelRCO.PRATICA),
        ("AULA_03_PRÁTICA_FRONT_END.docx",      PapelRCO.PRATICA),
        ("AULA_05_PRACTICE_WEB.docx",           PapelRCO.PRATICA),
        # Slides
        ("AULA_03_FRONT_END.pptx",              PapelRCO.SLIDES),
        ("AULA_03_PROGRAMACAO_WEB.pptx",        PapelRCO.SLIDES),
        ("aula_07_matematica.ppt",              PapelRCO.SLIDES),
        # Outro
        ("README.txt",                          PapelRCO.OUTRO),
        ("material_extra.pdf",                  PapelRCO.OUTRO),
    ])
    def test_detecta_papel_corretamente(self, nome, esperado):
        assert detectar_papel_rco(nome) == esperado

    def test_atividade_tem_prioridade_sobre_slides(self):
        """Arquivo DOCX com 'ATIVIDADE' não deve ser confundido com slides."""
        papel = detectar_papel_rco("AULA_03_ATIVIDADE_FRONT_END.docx")
        assert papel == PapelRCO.ATIVIDADE
        assert papel != PapelRCO.SLIDES

    def test_pratica_nunca_confundida_com_atividade(self):
        """PRATICA e ATIVIDADE devem ser sempre distintos — regra crítica."""
        assert detectar_papel_rco("AULA_03_PRATICA_X.docx")   == PapelRCO.PRATICA
        assert detectar_papel_rco("AULA_03_ATIVIDADE_X.docx") == PapelRCO.ATIVIDADE


class TestDetectarNumeroAula:

    @pytest.mark.parametrize("nome,esperado", [
        ("AULA_03_FRONT_END.pptx", 3),
        ("AULA_10_ATIVIDADE_WEB.docx", 10),
        ("aula_01_introducao.pptx", 1),
        ("material_sem_numero.pdf", None),
    ])
    def test_extrai_numero_corretamente(self, nome, esperado):
        assert detectar_numero_aula(nome) == esperado


# ── Testes: montar_conjunto_de_uploads ────────────────────────────────────────

class TestMontarConjuntoDeUploads:

    def _make_upload_file(self, nome: str):
        f = MagicMock()
        f.name = nome
        return f

    def test_monta_conjunto_com_3_arquivos(self):
        arquivos = {
            "arquivo_0": self._make_upload_file("AULA_03_FRONT_END.pptx"),
            "arquivo_1": self._make_upload_file("AULA_03_ATIVIDADE_FRONT_END.docx"),
            "arquivo_2": self._make_upload_file("AULA_03_PRATICA_FRONT_END.docx"),
        }
        conjunto = montar_conjunto_de_uploads(arquivos)

        assert conjunto.slides is not None
        assert conjunto.atividade is not None
        assert conjunto.pratica is not None
        assert conjunto.completo

    def test_conjunto_incompleto_sem_ppt(self):
        arquivos = {
            "arquivo_0": self._make_upload_file("AULA_03_ATIVIDADE_FRONT_END.docx"),
        }
        conjunto = montar_conjunto_de_uploads(arquivos)

        assert not conjunto.completo

    def test_pratica_opcional(self):
        """Conjunto com slides e atividade, sem prática, ainda é válido (completo=True)."""
        arquivos = {
            "arquivo_0": self._make_upload_file("AULA_03_FRONT_END.pptx"),
            "arquivo_1": self._make_upload_file("AULA_03_ATIVIDADE_FRONT_END.docx"),
        }
        conjunto = montar_conjunto_de_uploads(arquivos)

        assert conjunto.completo
        assert conjunto.pratica is None


# ── Testes: extrair_rco ───────────────────────────────────────────────────────

class TestExtrairRco:

    def test_extrai_todos_os_papeis(self):
        """extrair_rco() deve retornar dict com conteúdo para cada papel."""
        slides_buf    = _make_pptx_buffer("Aula 3 — CSS Flexbox")
        atividade_buf = _make_docx_buffer("1. O que é Flexbox?")
        pratica_buf   = _make_docx_buffer("Crie um layout com display: flex.")

        slides_buf.name    = "AULA_03_FRONT_END.pptx"
        atividade_buf.name = "AULA_03_ATIVIDADE_FRONT_END.docx"
        pratica_buf.name   = "AULA_03_PRATICA_FRONT_END.docx"

        conjunto = ConjuntoRCO(
            slides=slides_buf,
            atividade=atividade_buf,
            pratica=pratica_buf,
        )
        resultado = extrair_rco(conjunto)

        assert PapelRCO.SLIDES    in resultado
        assert PapelRCO.ATIVIDADE in resultado
        assert PapelRCO.PRATICA   in resultado

    def test_slides_tem_conteudo(self):
        """O papel 'slides' deve ter conteúdo extraído do PPTX."""
        slides_buf = _make_pptx_buffer("CSS Flexbox")
        slides_buf.name = "AULA_03_FRONT_END.pptx"

        conjunto = ConjuntoRCO(slides=slides_buf)
        resultado = extrair_rco(conjunto)

        assert resultado[PapelRCO.SLIDES].ok
        assert "CSS Flexbox" in resultado[PapelRCO.SLIDES].conteudo

    def test_atividade_tem_conteudo(self):
        """O papel 'atividade' deve ter conteúdo extraído do DOCX de questões."""
        atividade_buf = _make_docx_buffer("1. Questão de HTML")
        atividade_buf.name = "AULA_03_ATIVIDADE_FRONT_END.docx"

        conjunto = ConjuntoRCO(atividade=atividade_buf)
        resultado = extrair_rco(conjunto)

        assert resultado[PapelRCO.ATIVIDADE].ok
        assert "Questão de HTML" in resultado[PapelRCO.ATIVIDADE].conteudo

    def test_pratica_ausente_retorna_vazio(self):
        """Quando prática não é fornecida, seu conteúdo deve ser string vazia."""
        conjunto = ConjuntoRCO()  # sem arquivos
        resultado = extrair_rco(conjunto)

        assert resultado[PapelRCO.PRATICA].conteudo == ""


# ── Testes: montar_conteudo_rco ───────────────────────────────────────────────

class TestMontarConteudoRco:

    def test_retorna_dict_com_3_chaves(self):
        """montar_conteudo_rco() deve retornar dict com 'slides', 'atividade', 'pratica'."""
        extracao = {
            PapelRCO.SLIDES:    ResultadoExtracao("Conteúdo slides", TipoArquivo.PPTX),
            PapelRCO.ATIVIDADE: ResultadoExtracao("Questões ATIVIDADE", TipoArquivo.DOCX),
            PapelRCO.PRATICA:   ResultadoExtracao("Passos PRATICA", TipoArquivo.DOCX),
        }
        conteudo = montar_conteudo_rco(extracao)

        assert "slides"    in conteudo
        assert "atividade" in conteudo
        assert "pratica"   in conteudo

    def test_atividade_nao_e_pratica(self):
        """
        REGRA CRÍTICA: conteúdo da atividade NÃO pode ser o mesmo que a prática.
        ATIVIDADE → seção 12 (questões de fixação)
        PRATICA   → seção 11 (atividade prática)
        """
        extracao = {
            PapelRCO.SLIDES:    ResultadoExtracao("Slides", TipoArquivo.PPTX),
            PapelRCO.ATIVIDADE: ResultadoExtracao("CONTEUDO_ATIVIDADE_UNICO", TipoArquivo.DOCX),
            PapelRCO.PRATICA:   ResultadoExtracao("CONTEUDO_PRATICA_UNICO", TipoArquivo.DOCX),
        }
        conteudo = montar_conteudo_rco(extracao)

        assert conteudo["atividade"] == "CONTEUDO_ATIVIDADE_UNICO"
        assert conteudo["pratica"]   == "CONTEUDO_PRATICA_UNICO"
        # Garante que não foram invertidos
        assert conteudo["atividade"] != conteudo["pratica"]

    def test_pratica_vazia_quando_nao_fornecida(self):
        """Quando prática não existe, o campo deve ser string vazia."""
        extracao = {
            PapelRCO.SLIDES:    ResultadoExtracao("Slides", TipoArquivo.PPTX),
            PapelRCO.ATIVIDADE: ResultadoExtracao("Questões", TipoArquivo.DOCX),
            PapelRCO.PRATICA:   ResultadoExtracao("", TipoArquivo.DOCX),
        }
        conteudo = montar_conteudo_rco(extracao)

        assert conteudo["pratica"] == ""


# ── Teste de integração RCO completo ─────────────────────────────────────────

class TestPipelineRcoCompleto:

    def test_pipeline_completo_rco(self):
        """
        Teste de integração: upload → detectar papéis → extrair → montar conteúdo.
        Garante que cada papel chega no campo correto e não há inversão.
        """
        slides_buf    = _make_pptx_buffer("CSS Grid — Aula 05")
        atividade_buf = _make_docx_buffer("1. O que é CSS Grid?\nA) Sistema de layout\nB) Seletor")
        pratica_buf   = _make_docx_buffer("Crie uma grade de 3 colunas usando CSS Grid.")

        # Simula request.FILES
        slides_buf.name    = "AULA_05_FRONT_END.pptx"
        atividade_buf.name = "AULA_05_ATIVIDADE_FRONT_END.docx"
        pratica_buf.name   = "AULA_05_PRATICA_FRONT_END.docx"

        uploads = {
            "arquivo_0": slides_buf,
            "arquivo_1": atividade_buf,
            "arquivo_2": pratica_buf,
        }

        # Passo 1: montar conjunto
        conjunto = montar_conjunto_de_uploads(uploads)
        assert conjunto.completo

        # Passo 2: extrair
        extracao = extrair_rco(conjunto)

        # Passo 3: montar conteúdo final
        conteudo = montar_conteudo_rco(extracao)

        # Verifica slides
        assert "CSS Grid" in conteudo["slides"]

        # REGRA CRÍTICA: atividade contém questões, prática contém o exercício
        assert "O que é CSS Grid" in conteudo["atividade"]
        assert "Crie uma grade" in conteudo["pratica"]

        # Garante que os campos não foram invertidos
        assert "Crie uma grade" not in conteudo["atividade"]
        assert "O que é CSS Grid" not in conteudo["pratica"]
